"""Génère le support de soutenance (.pptx) de la plateforme MiamGo / FitMeal.

Deck jury (16 slides, français) : découpe microservices, architecture, technologies,
communication inter-services et l'ensemble des diagrammes de l'énoncé (rendus en
images depuis documentation/diagrammes.md via mermaid-cli).

Usage : python livrables/tools/build_presentation.py
Sortie : livrables/6_presentation_soutenance.pptx
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Chemins ---------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[2]
IMG = ROOT / "livrables" / "diagrammes_images"
OUT = ROOT / "livrables" / "6_presentation_soutenance.pptx"

# --- Design system ---------------------------------------------------------
DEEP = RGBColor(0x08, 0x42, 0x7B)      # bleu profond C4 (barre de titre, titre)
MID = RGBColor(0x11, 0x68, 0xBD)       # bleu C4 (accents secondaires)
ACCENT = RGBColor(0xE8, 0x59, 0x0C)    # orange MiamGo (liseré, puces)
INK = RGBColor(0x1F, 0x2A, 0x37)       # texte principal
MUTED = RGBColor(0x5B, 0x66, 0x72)     # texte secondaire
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)     # fond de zone claire
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xE7, 0xEE, 0xF6)      # entête de tableau
FONT = "Calibri"

EMU_IN = 914400
SW, SH = Inches(13.333), Inches(7.5)

TOTAL = 16


def solid(shape, color) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def set_run(run, *, size, color=INK, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_bar(slide, title, index):
    """Barre de titre colorée + liseré orange + index de section."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    solid(bar, DEEP)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SW, Inches(0.06))
    solid(stripe, ACCENT)

    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=26, color=WHITE, bold=True)

    idx = textbox(slide, SW - Inches(1.7), Inches(0.32), Inches(1.4), Inches(0.5))
    idx.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = idx.text_frame.paragraphs[0].add_run()
    r.text = f"{index:02d} / {TOTAL:02d}"
    set_run(r, size=12, color=RGBColor(0xBF, 0xD3, 0xEA), bold=True)


def footer(slide, page):
    ft = textbox(slide, Inches(0.5), SH - Inches(0.45), Inches(9), Inches(0.35))
    r = ft.text_frame.paragraphs[0].add_run()
    r.text = "MiamGo / FitMeal — Soutenance Architecture Microservices"
    set_run(r, size=9, color=MUTED)
    pg = textbox(slide, SW - Inches(1.3), SH - Inches(0.45), Inches(0.9), Inches(0.35))
    pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = pg.text_frame.paragraphs[0].add_run()
    r.text = str(page)
    set_run(r, size=9, color=MUTED, bold=True)


def content_slide(prs, title, index):
    s = blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    solid(bg, WHITE)
    bg.shadow.inherit = False
    # renvoyer au fond
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    title_bar(s, title, index)
    footer(s, index)
    return s


def add_bullets(slide, items, *, left=Inches(0.7), top=Inches(1.5),
                width=Inches(12.0), height=Inches(5.3), size=18):
    """items: list of (text, level, kind) ; kind in {'bullet','head','note'}."""
    tb = textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, level, kind = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        p.space_before = Pt(2)
        if kind == "head":
            r = p.add_run()
            r.text = text
            set_run(r, size=size + 3, color=DEEP, bold=True)
        elif kind == "note":
            r = p.add_run()
            r.text = text
            set_run(r, size=size - 3, color=MUTED, italic=True)
        else:
            bullet = "▸ " if level == 0 else "– "
            rb = p.add_run()
            rb.text = bullet
            set_run(rb, size=size, color=ACCENT, bold=True)
            r = p.add_run()
            r.text = text
            set_run(r, size=size - (2 if level else 0), color=INK)
    return tb


def add_image_fit(slide, name, *, top=Inches(1.45), bottom_margin=Inches(0.55),
                  side_margin=Inches(0.5), caption=None):
    path = IMG / name
    with Image.open(path) as im:
        iw, ih = im.size
    avail_w = SW - side_margin * 2
    cap_h = Inches(0.4) if caption else Inches(0)
    avail_h = SH - top - bottom_margin - cap_h
    scale = min(avail_w / iw, avail_h / ih)
    w = Emu(int(iw * scale))
    h = Emu(int(ih * scale))
    left = Emu(int((SW - w) / 2))
    slide.shapes.add_picture(str(path), left, top, width=w, height=h)
    if caption:
        cb = textbox(slide, Inches(0.5), top + h + Inches(0.05), SW - Inches(1.0), Inches(0.4))
        cb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = cb.text_frame.paragraphs[0].add_run()
        r.text = caption
        set_run(r, size=11, color=MUTED, italic=True)


def add_table(slide, headers, rows, *, left=Inches(0.6), top=Inches(1.5),
              width=Inches(12.1), col_widths=None, header_size=13, body_size=12,
              row_height=Inches(0.42)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_height * n_rows
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gtable.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    # entête
    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htxt
        set_run(r, size=header_size, color=WHITE, bold=True)
    # corps
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(6)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = val
            bold = j == 0
            set_run(r, size=body_size, color=DEEP if bold else INK, bold=bold)
    return table


# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ---- Slide 1 : Titre --------------------------------------------------
    s = blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    solid(bg, DEEP)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.05), SW, Inches(0.08))
    solid(band, ACCENT)
    t = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6))
    p = t.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "MiamGo / FitMeal"
    set_run(r, size=54, color=WHITE, bold=True)
    p2 = t.text_frame.add_paragraph()
    r = p2.add_run(); r.text = "Architecture microservices d'une plateforme de livraison de repas"
    set_run(r, size=24, color=RGBColor(0xBF, 0xD3, 0xEA))
    sub = textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.2))
    for i, line in enumerate([
        ("Conception, prototype fullstack et patterns de résilience", 18, WHITE, True),
        ("FastAPI · API Gateway Nginx · Redis Pub/Sub · SAGA orchestrée · Circuit Breaker · React/TypeScript", 14, RGBColor(0xBF, 0xD3, 0xEA), False),
        ("Soutenance de projet — Architecture Microservices · 2026", 13, RGBColor(0x9D, 0xB8, 0xD6), False),
    ]):
        p = sub.text_frame.paragraphs[0] if i == 0 else sub.text_frame.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = line[0]
        set_run(r, size=line[1], color=line[2], bold=line[3])

    # ---- Slide 2 : Contexte & problématique -------------------------------
    s = content_slide(prs, "Contexte métier & problématique", 2)
    add_bullets(s, [
        ("Trois acteurs aux besoins divergents", 0, "head"),
        ("Client : recherche rapide, commande fluide, suivi en temps réel", 1, "bullet"),
        ("Restaurateur : gestion de la carte, validation et cadence de cuisine", 1, "bullet"),
        ("Livreur : affectation géolocalisée et confirmation de course", 1, "bullet"),
        ("Défis techniques structurants", 0, "head"),
        ("Forte asymétrie de charge : consultation intensive du catalogue vs transactions d'achat critiques", 1, "bullet"),
        ("Transactions distribuées : pas de BDD unique → cohérence entre services sans verrou ACID", 1, "bullet"),
        ("Résilience : pannes possibles des dépendances externes (PSP de paiement, réseau)", 1, "bullet"),
        ("→ Une architecture microservices découpée par charge, événementielle et résiliente", 0, "note"),
    ])

    # ---- Slide 3 : Découpe microservices (DDD) ----------------------------
    s = content_slide(prs, "Découpe en microservices (DDD, par charge)", 3)
    add_bullets(s, [("6 microservices autonomes + API Gateway — découpe par charge et par bounded context (ADR 0001 & 0002)", 0, "note")],
                top=Inches(1.35), height=Inches(0.5), size=15)
    add_table(s,
              ["Microservice", "Bounded context", "Motif de découpe"],
              [
                  ["service-utilisateurs", "Identité & comptes", "Source de vérité RBAC (client / restaurant_owner / courier)"],
                  ["service-restaurants", "Catalogue & menus", "Forte consultation, isolée des écritures de commande"],
                  ["service-commandes", "Panier & checkout", "Orchestrateur SAGA et machine à états de la commande"],
                  ["service-paiements", "Paiement & PSP", "Criticité maximale : idempotence + circuit breaker"],
                  ["service-livraisons", "Flotte & tracking", "Attribution géolocalisée (haversine) des livreurs"],
                  ["service-notifications", "Notifications", "Consommateur d'événements, envois simulés (email/push/sms)"],
              ],
              top=Inches(1.95),
              col_widths=[Inches(2.7), Inches(2.6), Inches(6.8)],
              row_height=Inches(0.62), body_size=12.5)

    # ---- Slide 4 : C4 niveau 1 -------------------------------------------
    s = content_slide(prs, "Architecture — Contexte système (C4 niveau 1)", 4)
    add_image_fit(s, "c4-contexte.png",
                  caption="Le système, ses trois acteurs et le PSP externe simulé")

    # ---- Slide 5 : C4 niveau 2 -------------------------------------------
    s = content_slide(prs, "Architecture — Conteneurs (C4 niveau 2)", 5)
    add_image_fit(s, "c4-conteneurs.png",
                  caption="Gateway Nginx, 6 microservices FastAPI, bus Redis et stores isolés par service")

    # ---- Slide 6 : Technologies ------------------------------------------
    s = content_slide(prs, "Technologies utilisées", 6)
    add_table(s,
              ["Couche", "Technologies"],
              [
                  ["Backend", "Python 3.13 · FastAPI 0.139 · Pydantic 2 · Uvicorn · httpx · structlog (logs JSON)"],
                  ["Communication", "REST/HTTP (httpx, synchrone) · Redis 7 Pub/Sub (événements asynchrones)"],
                  ["Frontend", "React 18 · TypeScript · Vite · TailwindCSS · shadcn/ui (Radix) · TanStack Query · axios"],
                  ["Infrastructure", "Docker Compose · API Gateway Nginx · Redis 7-alpine · healthchecks par service"],
                  ["Qualité", "pytest + coverage (>95%) · ruff · mypy (strict) · Vitest · ESLint · Prettier"],
              ],
              top=Inches(1.9),
              col_widths=[Inches(2.4), Inches(9.7)],
              row_height=Inches(0.72), body_size=13)
    add_bullets(s, [("Stack homogène par service · configuration par variables d'environnement préfixées (pydantic-settings)", 0, "note")],
                top=Inches(5.9), height=Inches(0.5), size=13)

    # ---- Slide 7 : Communication inter-services --------------------------
    s = content_slide(prs, "Communication inter-services", 7)
    add_bullets(s, [
        ("Synchrone (REST/HTTP via httpx) — phase critique de la SAGA", 0, "head"),
        ("service-commandes est le seul appelant sortant : validation, paiement, ticket cuisine, livraison", 1, "bullet"),
        ("Asynchrone (Redis Pub/Sub) — découplage et continuation événementielle", 0, "head"),
    ], top=Inches(1.4), height=Inches(2.0), size=15)
    add_table(s,
              ["Canal d'événement", "Producteur", "Consommateur(s)"],
              [
                  ["order.confirmed / cancelled / delivered", "commandes", "notifications"],
                  ["order.ready", "restaurants", "commandes (→ attribution livreur)"],
                  ["delivery.assigned / picked_up", "livraisons", "notifications"],
                  ["delivery.completed", "livraisons", "commandes (→ DELIVERED) + notifications"],
              ],
              top=Inches(3.5),
              col_widths=[Inches(4.6), Inches(2.4), Inches(5.1)],
              row_height=Inches(0.5), body_size=12.5)

    # ---- Slide 8 : SAGA vue d'ensemble -----------------------------------
    s = content_slide(prs, "Transaction distribuée : SAGA orchestrée", 8)
    add_bullets(s, [
        ("Orchestration hybride (ADR 0003) : phase critique synchrone, puis continuation chorégraphiée", 0, "note"),
        ("Étape 1 — Création de la commande (statut RECEIVED)", 0, "bullet"),
        ("Étape 2 — Validation restaurant (sous-total, disponibilité des plats)", 0, "bullet"),
        ("Étape 3 — Paiement sécurisé (circuit breaker + retry + timeout)", 0, "bullet"),
        ("Étape 4 — Ticket cuisine (acceptation par le restaurant)", 0, "bullet"),
        ("Étape 5 — Confirmation (statut PREPARING) + événement order.confirmed", 0, "bullet"),
        ("Transactions compensatoires", 0, "head"),
        ("Refus cuisine après paiement capturé → remboursement intégral automatique (CANCELLED_REFUSED)", 1, "bullet"),
        ("Aucun livreur disponible → compensation tardive : remboursement + annulation", 1, "bullet"),
        ("État de la SAGA persisté à chaque étape (saga_state) → progression observable ; la SAGA ne lève jamais d'erreur HTTP", 0, "note"),
    ], top=Inches(1.4), size=16)

    # ---- Slide 9 : SAGA séquence nominale --------------------------------
    s = content_slide(prs, "SAGA — séquence nominale du passage de commande", 9)
    add_image_fit(s, "saga-nominale.png",
                  caption="Checkout réussi : validation → paiement → ticket cuisine → confirmation")

    # ---- Slide 10 : Résilience patterns ----------------------------------
    s = content_slide(prs, "Pattern de résilience (ADR 0007)", 10)
    add_bullets(s, [
        ("Trois patterns combinés sur les appels sortants de la SAGA", 0, "head"),
        ("Timeout : 2,0 s sur chaque appel HTTP amont", 0, "bullet"),
        ("Retry : 3 tentatives, backoff exponentiel + jitter — uniquement sur erreurs transitoires (5xx, réseau), jamais sur un 4xx", 0, "bullet"),
        ("Circuit Breaker (CLOSED → OPEN → HALF_OPEN) : ouverture à 5 échecs / 30 s, appliqué sur l'appel paiement", 0, "bullet"),
        ("Sûreté & démonstration", 0, "head"),
        ("Idempotence par order_id côté paiements → le retry ne provoque jamais de double débit", 0, "bullet"),
        ("PSP Chaos Mode : endpoint POST /_chaos règle le taux d'échec à chaud pour la démo QA", 0, "bullet"),
        ("Circuit ouvert → annulation propre de la commande, sans débit client", 0, "note"),
    ], top=Inches(1.4), size=17)

    # ---- Slide 11 : Résilience séquence ----------------------------------
    s = content_slide(prs, "Résilience — panne PSP & Circuit Breaker OPEN", 11)
    add_image_fit(s, "saga-resilience.png",
                  caption="Retries avec backoff, ouverture du circuit, annulation sans débit")

    # ---- Slide 12 : Compensation -----------------------------------------
    s = content_slide(prs, "SAGA — compensation (remboursement automatique)", 12)
    add_image_fit(s, "saga-compensation.png",
                  caption="Refus cuisine après capture → remboursement intégral via le PSP")

    # ---- Slide 13 : Chorégraphie livraison -------------------------------
    s = content_slide(prs, "Continuation livraison (chorégraphie événementielle)", 13)
    add_image_fit(s, "choregraphie-livraison.png",
                  caption="order.ready → attribution livreur → suivi → delivery.completed → DELIVERED")

    # ---- Slide 14 : Données & isolation ----------------------------------
    s = content_slide(prs, "Gestion des données & isolation (ADR 0005)", 14)
    add_image_fit(s, "erd-donnees.png", top=Inches(1.35), bottom_margin=Inches(0.5),
                  caption="Un modèle et un store autonomes par service — liens logiques par clés d'agrégat")

    # ---- Slide 15 : Frontend & rôles -------------------------------------
    s = content_slide(prs, "Application fullstack & rôles (RBAC)", 15)
    add_bullets(s, [
        ("Frontend React 18 / TypeScript / Vite / Tailwind — 4 vues cloisonnées par rôle (ADR 0010)", 0, "note"),
        ("Vue Client : catalogue, recherche filtrée, panier, checkout GPS, suivi en direct", 0, "bullet"),
        ("Espace Restaurateur : établissement, gestion de la carte, tickets cuisine (PREPARING → READY)", 0, "bullet"),
        ("Espace Livreur : flotte géolocalisée, prise en charge et remise (PICKED_UP → DELIVERED)", 0, "bullet"),
        ("Dashboard QA Testeur : santé des 6 services, PSP Chaos, connexion instantanée aux comptes de test", 0, "bullet"),
        ("Contrôle d'accès par rôle : garde RequireRole côté routes, rôle porté par le service-utilisateurs", 0, "note"),
    ], top=Inches(1.5), size=17)

    # ---- Slide 16 : Démo / Qualité / Perspectives / Q&R ------------------
    s = content_slide(prs, "Démonstration · Qualité · Perspectives", 16)
    add_bullets(s, [
        ("Scénarios de démonstration", 0, "head"),
        ("Parcours nominal : commande → PREPARING → READY → attribution livreur → DELIVERED", 1, "bullet"),
        ("Résilience : PSP Chaos 100 % → circuit breaker → annulation « sans débit »", 1, "bullet"),
        ("Qualité", 0, "head"),
        ("Backend : > 95 % de couverture (~310 tests), ruff + mypy strict ; Frontend : Vitest, ESLint, Prettier", 1, "bullet"),
        ("Conteneurisation Docker Compose avec healthchecks sur tous les services", 1, "bullet"),
        ("Perspectives", 0, "head"),
        ("Kafka (persistance/replay) · CQRS lecture catalogue · PostgreSQL + Alembic · JWT RS256 au gateway", 1, "bullet"),
        ("Merci — questions & réponses", 0, "note"),
    ], top=Inches(1.4), size=16)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK -> {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
